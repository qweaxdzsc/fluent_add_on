
project_name = 'D2U-2'
version = 'V20.2_vent'
Rotation_speed = 3600


def get_ppt(project_name, version, RPM):
    """organize data from database to pptx
    how is works:
    1. find corresponding data from mysql
    2. create ppt style for each page
    """

    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    import time
    import os
    import numpy as np
    from pptx.enum.text import MSO_ANCHOR
    from pptx.enum.text import PP_ALIGN
    import pymysql

    # extract info from mysql
    user_name = 'root'
    pwd = 'Sdaac=12345'

    conn = pymysql.connect(host='localhost',
                           user=user_name,
                           password=pwd,
                           db='cfd-result',
                           charset='utf8mb4'
                           )


    # search main table
    cursor = conn.cursor(pymysql.cursors.DictCursor)  # return Dictionary_like key_valve
    serh = """
    SELECT  * FROM cfd_project where Project = '%s' and Version = '%s' and RPM = '%s'
    """ % (project_name, version, RPM)
    cursor.execute(serh)
    main_data = cursor.fetchone()
    print('main table data:', main_data)
    cursor.close()

    # search volume table
    ID = main_data['ID']
    cursor = conn.cursor()
    s_volume = """
    SELECT `Face_name`, `Volume(L/S)`, `Percentage` FROM cfd_volume where Project_id = '%s'
    """ % (ID)
    cursor.execute(s_volume)
    volume_data = cursor.fetchall()
    print(volume_data)
    cursor.close()

    # search sp table
    cursor = conn.cursor()
    s_sp = """
        SELECT `Face_name`, `Static_pressure` FROM cfd_sp where Project_id = '%s'
        """ % (ID)
    cursor.execute(s_sp)
    sp_data = cursor.fetchall()
    print(sp_data)
    cursor.close()

    # search tp table
    cursor = conn.cursor()
    s_tp = """
        SELECT `Face_name`, `Total_pressure` FROM cfd_tp where Project_id = '%s'
        """ % (ID)
    cursor.execute(s_tp)
    tp_data = cursor.fetchall()
    print(tp_data)
    cursor.close()

    # search uni table
    cursor = conn.cursor()
    s_uni = """
            SELECT `Uniformity` FROM cfd_uni where Project_id = '%s' and Face_name = '%s'
            """ % (ID, 'evap_out')
    cursor.execute(s_uni)
    uni_data = cursor.fetchone()
    print('This is uni data:', uni_data)
    cursor.close()
    conn.close()

    # info pre-process
    result_path = main_data['File_address']
    RPM = str(main_data['RPM'])
    total_volume = str(main_data['Total_volume'])
    volume_data = [list(x) for x in volume_data]
    torque = str(main_data['Torque'])
    fan_efficiency = '%.1f %%'% (main_data['Fan_efficiency']*100)

    whole_title = project_name + '-' + version
    evap_out_uni = str(uni_data[0])

    # start PPT writing
    ppt = Presentation(r'C:\Users\BZMBN4\Documents\Custom Office Templates\template.pptx')  # call template
    # ppt = Presentation(r'C:\Users\BZMBN4\Desktop\template.pptx')

    def cover_slide(title):
        cover_slide = ppt.slides[0]
        title_shape = cover_slide.shapes.placeholders[0]
        title_shape.text_frame.clear()
        p = title_shape.text_frame.paragraphs[0]
        run1 = p.add_run()
        run1.text = '%s\n\n' % (title)
        run1.font.name = 'Arial Unicode MS'
        run1.font.size = Pt(30)
        # run1.font.italic = True
        run1.font.color.rgb = RGBColor(95, 99, 102)

        run2 = p.add_run()
        run2.text = time.strftime('%Y/%m/%d', time.localtime(time.time()))
        run2.font.name = 'Arial Unicode MS'
        run2.font.size = Pt(18)
        # run2.font.italic = True
        run2.font.color.rgb = RGBColor(49, 163, 181)

    def title_set(slide, title_name):
        title = slide.shapes[0].text_frame.paragraphs[0]
        title.text = '%s' %(title_name)
        title.font.name = 'Arial Unicode MS'
        title.font.size = Pt(28)
        title.font.color.rgb = RGBColor(46, 117, 182)

    def place_table(row, col, left, top, width, height, slide):

        table = slide.shapes.add_table(row, col, left, top, width, height).table     # create table
        for j in range(row):                                                         # change color for table
            if j == 0:
                for i in range(col):
                    table.cell(0, i).fill.solid()
                    table.cell(0, i).fill.fore_color.rgb = RGBColor(68, 114, 196)
            elif j % 2 == 0 and j > 0:
                for i in range(col):
                    table.cell(j, i).fill.solid()
                    table.cell(j, i).fill.fore_color.rgb = RGBColor(233, 235, 245)
            else:
                for i in range(col):
                    table.cell(j, i).fill.solid()
                    table.cell(j, i).fill.fore_color.rgb = RGBColor(207, 213, 234)
        return table

    # create cover slide
    cover_slide(project_name)

    # create overview slide
    overview_slide = ppt.slides[1]
    title_set(overview_slide, project_name+'-' + version + '-CFD-Result')

    col = 7                                             # create table
    row = 2
    left = Inches((1+(7-col)*0.4)/2)
    top = Inches(2-row*0.1)
    width = Inches(5-(7-col)*0.4)
    height = Inches(0.5*row)
    ov_table = place_table(row, col, Inches(0.15), Inches(1.6), width, height, overview_slide)

    # input data
    heading = ['Project', 'Version', 'RPM', 'Volume\n(L/S)', 'Uniformity\nEvap_out', 'Torque\n(N*M)', 'Fan\nefficiency']
    for i in range(col):
        ov_table.cell(0, i).text = heading[i]

    ov_table.cell(1, 0).text = project_name
    ov_table.cell(1, 1).text = version
    ov_table.cell(1, 2).text = RPM
    ov_table.cell(1, 3).text = total_volume
    ov_table.cell(1, 4).text = evap_out_uni
    ov_table.cell(1, 5).text = torque
    ov_table.cell(1, 6).text = fan_efficiency

    # aligning center the table text
    def aligning_table(table):
        row = len(table.rows)
        col = len(table.columns)
        lenth_matrix = np.zeros([row, col])
        for i in range(row):
            for j in range(col):
                table.cell(i, j).vertical_anchor = MSO_ANCHOR.MIDDLE                     # vertical center alignment
                for k in range(len(table.cell(i, j).text_frame.paragraphs)):
                    table.cell(i, j).text_frame.paragraphs[k].alignment = PP_ALIGN.CENTER    # horizontal center alignment
                lenth_matrix[i, j] = len(table.cell(i, j).text)

        return lenth_matrix

    def auto_wid_table(col, lenth_matrix, table):
        str_wid = np.zeros(col)
        total_wid = 0
        for i in range(col):                                                                   # table auto-width
            str_wid[i] = max(lenth_matrix[:, i])*2.5 + 46
            table.columns[i].width = Pt(str_wid[i])
            total_wid += str_wid[i]

    def table_font(table, head_size, body_size):
        for i in range(row):
            if i == 0:
                for j in range(col):
                    table.cell(0, j).text_frame.paragraphs[0].font.size = head_size
            else:
                for j in range(col):
                    table.cell(i, j).text_frame.paragraphs[0].font.size = body_size
                    table.cell(i, j).text_frame.paragraphs[0].font.bold = True

    word_lenth = aligning_table(ov_table)
    auto_wid_table(col, word_lenth, ov_table)
    table_font(ov_table, Pt(12), Pt(12))

    # create distribution slide
    dis_slide = ppt.slides[2]
    print(len(dis_slide.shapes))
    title_set(dis_slide, whole_title + '-distribution')

    col = 3                                                     # place table
    row = len(volume_data) + 1
    left = Inches((0.7+(7-col)*0.4)/2)
    top = Inches(2.4-row*0.1)
    height = Inches(0.4*row)
    dis_table = place_table(row, col, left, top, Inches(5), height, dis_slide)

    # put data in
    dis_table_head = ['', 'Volume(L/S)', 'Percentage(%)']
    for i in range(len(dis_table_head)):
        dis_table.cell(0, i).text = dis_table_head[i]

    for i in range(row-1):
        for j in range(len(volume_data[0])):
            if j == 2:
                volume_data[i][j] = '%.1f%%' % (((volume_data[i][j]))*100)
            dis_table.cell(i+1, j).text = str(volume_data[i][j])

    # modify font in table
    table_font(dis_table, Pt(12), Pt(11))
    aligning_table(dis_table)

    def set_picture(result_path, slide, picture):
        file_path = result_path + '\\'

        left = Inches(0.8)
        top = Inches(1.3)
        width = Inches(6)
        slide.shapes.add_picture(file_path + picture, left, top, width)

    def dp_table_create(table, head, dp_data):
        for j in range(col):
            table.cell(0, j).text = head[j]

        for i in range(row-1):
            for j in range(len(dp_data[0])):
                table.cell(i+1, j).text = str(dp_data[i][j])

        table_font(table, Pt(10), Pt(10))
        aligning_table(table)

    # Pressure Drop slide
    dp_slide = ppt.slides[3]
    title_set(dp_slide, whole_title + '-Pressure-Drop')

    # create table
    col = 2
    row = len(sp_data) + 1
    left = Inches(0.8)
    top = Inches(1.2)
    width = Inches(2.5)
    height = Inches(0.35 * row)

    sp_table = place_table(row, col, left, top, width, height, dp_slide)
    sp_table_head = ['Static Pressure', version]
    dp_table_create(sp_table, sp_table_head, sp_data)

    tp_table = place_table(row, col, Inches(4), top, width, height, dp_slide)
    tp_table_head = ['Total Pressure', version]
    dp_table_create(tp_table, tp_table_head, tp_data)

    # picture slides
    picture_list = list()
    pic_name_list = list()
    total_picture_div = str()
    for file in os.listdir(result_path):
        if file.endswith('jpg') or file.endswith('png'):
            pic_name_list.append(file.split('.')[0])
            picture_list.append(file)
    print(pic_name_list)
    print(picture_list)
    for index, pic in enumerate(picture_list):
        if index < 3:
            wstrm_slide = ppt.slides[4 + index]
            title_set(wstrm_slide, '%s' % pic_name_list[index])
            set_picture(result_path, wstrm_slide, '%s' % pic)

    # save ppt and open it
    ppt.save(result_path + '\\' + whole_title + '.pptx')
    os.system(result_path + '\\' + whole_title + '.pptx')


if __name__ == "__main__":
    get_ppt(project_name, version, Rotation_speed)