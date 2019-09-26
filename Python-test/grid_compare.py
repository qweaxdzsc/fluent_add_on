import numpy as np

file_path = r'G:\GE2_REAR\GE2-rear-round2\GE2-rear-V20-FC\result_GE2-rear2_V20-FC'
file_name = file_path + '\\clip_velocity.txt'
f = open(file_name, encoding='utf-8')  # open
txt = f.readlines()
# read
ppt_title = 'GE2-Rear2-V20-FC-Grid-Uniformity'

for i in range(len(txt)):  # process txt
    txt[i] = txt[i].strip()  # strip blanks before and after
    txt[i] = txt[i].split()  # split line into list


def partitioning(txt):

    exist_factor = 0                       # determine the existence of part
    divide_sign = []

    for line in txt:
        if 'Velocity' and 'Magnitude' in line:               # find part
            title_index = txt.index(line)
            exist_factor += 1               # 1 means exist
            break
    if exist_factor == 1:                   # if can be found
        txt_down = txt[title_index:]        # have all line below title
        for line in txt_down:
            if '--' in line[0]:             # use part_name and "--" to isolate part
                divide_sign.append(txt_down.index(line))          # get index of top and bottom line
            if len(divide_sign) > 1:                              # find bottom, break
                break
                # print(divide_sign)
        txt_taget = txt_down[divide_sign[0]+1:divide_sign[1]]     # isolate part
        txt_taget = np.array(txt_taget)                           # convert to np for better manipulation
    else:
        txt_taget = ['not-exist']

    return txt_taget, exist_factor


txt_taget, exist_factor = partitioning(txt)
taget_data = txt_taget[:, 1]

data_array = np.array([float(i) for i in taget_data])
avg_grid = data_array.reshape([4, 4])
print('Average velocity grid:\n', avg_grid)

inters_box = np.zeros([7, 7])
inters_box[::2, ::2] = avg_grid[:, :]


def compare_grid(avg_grid):
    message = """"""
    show_box = np.zeros([4, 4])                         # for visualize result

    for i in range(len(avg_grid[1, :])):
        for j in range(len(avg_grid[:, 1])):

            # horizontal compare
            if j < 3:
                dif = avg_grid[i, j+1]-avg_grid[i, j]

                if dif >= 0:
                    perct = dif/avg_grid[i, j]
                else:
                    perct = abs(dif/avg_grid[i, j+1])

                if perct > 0.2:
                    show_box[i, j] += 1
                    show_box[i, j+1] += 1
                    show_per = perct*100
                    inters_box[i*2, j*2+1] += -1
                    message += 'WARNING: zone_index:%s-%s; zone_index:%s-%s. Difference:%.1f%%\n'\
                          %(i+1, j+1, i+1, j+2, show_per)
                    print('WARNING: zone_index:%s-%s, speed:%s; zone_index:%s-%s, speed:%s. Difference:%.1f%% '\
                          %(i+1, j+1, avg_grid[i, j], i+1, j+2, avg_grid[i, j+1], show_per))
            # vertical compare
            if i < 3:
                dif = avg_grid[i+1, j]-avg_grid[i, j]

                if dif >=0:
                    perct = dif/avg_grid[i, j]
                else:
                    perct = abs(dif/avg_grid[i+1, j])

                if perct > 0.2:
                    show_box[i, j] += 1
                    show_box[i+1, j] += 1
                    show_per = perct * 100
                    inters_box[i * 2 + 1, j * 2 ] += -3
                    message += 'WARNING: zone_index:%s-%s; zone_index:%s-%s. Difference:%.1f%%\n'\
                          %(i+1, j+1, i+2, j+1, show_per)
                    print('WARNING: zone_index:%s-%s, speed:%s; zone_index:%s-%s,speed:%s. Difference:%.1f%% '\
                          %(i+1, j+1, avg_grid[i, j], i+2, j+1, avg_grid[i+1, j], show_per))
            # prime-diag compare
            if i < 3 and j < 3:
                dif = avg_grid[i+1,j+1]-avg_grid[i, j]

                if dif >=0:
                    perct = dif/avg_grid[i, j]
                else:
                    perct = abs(dif/avg_grid[i+1, j+1])

                if perct > 0.2:
                    show_box[i, j] += 1
                    show_box[i+1, j+1] += 1
                    show_per = perct * 100
                    inters_box[i * 2 + 1, j * 2 + 1] += -2
                    message += 'WARNING: zone_index:%s-%s; zone_index:%s-%s. Difference:%.1f%%\n'\
                          %(i+1, j+1, i+2, j+2, show_per)
                    print('WARNING: zone_index:%s-%s, speed:%s; zone_index:%s-%s,speed:%s. Difference:%.1f%% '\
                          %(i+1, j+1, avg_grid[i, j], i+2, j+2, avg_grid[i+1, j+1], show_per))

            # vice-diag compare
            if i < 3 and j > 0:
                dif = avg_grid[i+1,j-1]-avg_grid[i, j]

                if dif >=0:
                    perct = dif/avg_grid[i, j]
                else:
                    perct = abs(dif/avg_grid[i+1, j-1])

                if perct > 0.2:
                    show_box[i, j] += 1
                    show_box[i+1, j-1] += 1
                    show_per = perct * 100
                    inters_box[i * 2 + 1, j * 2 - 1] += -4
                    message += 'WARNING: zone_index:%s-%s; zone_index:%s-%s. Difference:%.1f%%\n'\
                          %(i+1, j+1, i+2, j, show_per)
                    print('WARNING: zone_index:%s-%s, speed:%s; zone_index:%s-%s,speed:%s. Difference:%.1f%%  '\
                          %(i+1, j+1, avg_grid[i, j], i+2, j, avg_grid[i+1, j-1], show_per))

    return message, show_box, inters_box


message, show_box, inters_box = compare_grid(avg_grid)
print('Visualization(problem area represented by 1): \n', show_box)
print(inters_box)
f.close()

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

import os

# output to ppt
ppt = Presentation(r'C:\Users\BZMBN4\Documents\Custom Office Templates\evap_clip.pptx')
slide = ppt.slides[0]


# set title
def title_set(slide, title_name):
    title = slide.shapes[0].text_frame.paragraphs[0]
    title.text = '%s' %(title_name)
    title.font.name = 'Arial Unicode MS'
    title.font.size = Pt(28)
    title.font.color.rgb = RGBColor(46, 117, 182)


title_set(slide, ppt_title)

# add picture
img_path = file_path + '\\evap_clip.jpg'

ml = Inches(0.8)
mt = Inches(1.1)
width = Inches(8.5)
pic = slide.shapes.add_picture(img_path, ml, mt, width)

# add auto-shapes
all_shape = slide.shapes
star_left = Inches(2.3)
star_top = Inches(2)
star_width = star_height = Inches(0.3)

tbox_left = Inches(2.6)
tbox_top = Inches(1.8)
tbox_width = tbox_height = Inches(0.3)

star = [[0 for i in range(4)] for j in range(4)]
text_boxs = [[0 for i in range(4)] for j in range(4)]

for i in range(4):
    for j in range(4):
        text_boxs[i][j] = all_shape.add_textbox(tbox_left + j*Inches(1.42), tbox_top + i*Inches(1.08), tbox_width, tbox_height)
        text_boxs[i][j].text = str('%.3f' % (avg_grid[i, j]))
        text_boxs[i][j].text_frame.paragraphs[0].font.size = Pt(14)
        if show_box[i, j] == 0:
            star[i][j] = all_shape.add_shape(MSO_SHAPE.STAR_5_POINT, star_left + j*Inches(1.42), star_top + i*Inches(1.08), star_width, star_height)

arrow = [[0 for i in range(7)] for j in range(7)]
arrow_left = Inches(2.7)
arrow_top = Inches(1.8)
arrow_width = Inches(0.4)
arrow_height = Inches(0.2)
for i in range(7):
    for j in range(7):
        if inters_box[i, j] < 0:
            arrow[i][j] = all_shape.add_shape(MSO_SHAPE.LEFT_RIGHT_ARROW, arrow_left + j*Inches(0.715), arrow_top + i*Inches(0.55), arrow_width, arrow_height)
            arrow[i][j].rotation = abs(inters_box[i, j]+1)*45
            arrow[i][j].fill.solid()
            arrow[i][j].fill.fore_color.rgb = RGBColor(255, 0, 0)

note_left = Inches(2)
note_top = Inches(6.2)
note_width = Inches(8)
note_height = Inches(1)

note_boxs = all_shape.add_textbox(note_left, note_top, note_width, note_height)
note_boxs.text = message

note_para = note_boxs.text_frame.paragraphs
for i in range(len(note_para)):
    note_para[i].font.size = Pt(16)

ppt_name = file_path + '\\' + ppt_title + '.pptx'
ppt.save(ppt_name)
os.system(ppt_name)